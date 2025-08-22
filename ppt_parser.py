from pptx import Presentation
from pathlib import Path

def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        text_items = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_items.append(text)

        if text_items:
            slide_text = "\n".join(text_items)
            print(f"slide_text: {slide_text}")
            slides_text.append(f"📄 Slide {i} ({ppt_path.name}):\n{slide_text}")

    return "\n\n".join(slides_text)

def load_all_ppts(folder="pptx"):
    all_text = []

    for path in Path(folder).glob("*.pptx"):
        try:
            print(f"🖼️ Loading PowerPoint: {path.name}")
            text = extract_text_from_ppt(path)
            if text.strip():
                all_text.append({
                    "filename": path.name,
                    "content": text.strip()
                })
        except Exception as e:
            print(f"❌ Failed to parse {path.name}: {e}")

    return all_text
